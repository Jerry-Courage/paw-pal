import io
import logging
from typing import Dict, Any

logger = logging.getLogger('nitemind')

def extract_text_from_bytes(file_bytes: bytes, extension: str) -> Dict[str, Any]:
    """
    Extract text from common document formats.
    Returns: {'text': str, 'status': 'success'|'error', 'error': str}
    """
    content = {'text': '', 'status': 'success'}
    
    try:
        ext = extension.lower()
        if ext == '.pdf':
            from .pdf_extractor import extract_pdf_content
            pdf_data = extract_pdf_content(file_bytes=file_bytes)
            content['text'] = pdf_data['text']
            content['pdf_data'] = pdf_data # Keep images etc for caller
        
        elif ext in ['.docx', '.doc']:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            content['text'] = '\n'.join(full_text)
            
        elif ext in ['.txt', '.md', '.py', '.js', '.ts', '.css', '.html']:
            content['text'] = file_bytes.decode('utf-8', errors='ignore')

        elif ext in ['.pptx', '.ppt']:
            try:
                from pptx import Presentation
                import io as _io

                # ── Strategy 1: Convert to PDF via LibreOffice, then use the
                # full PDF extractor pipeline (identical quality to PDF uploads)
                pdf_bytes = _convert_pptx_to_pdf(file_bytes)
                if pdf_bytes:
                    logger.info("[PPTX] LibreOffice conversion succeeded — running full PDF pipeline")
                    from .pdf_extractor import extract_pdf_content
                    pdf_data = extract_pdf_content(file_bytes=pdf_bytes)
                    content['text'] = pdf_data['text']
                    content['pdf_data'] = pdf_data   # reuse PDF pipeline: page_images, images, toc
                    # Override ext so tasks.py treats it as PDF
                    content['converted_from_pptx'] = True
                    content['page_count'] = pdf_data.get('page_count', 0)

                    # Also extract text + speaker notes from PPTX for extra context
                    prs = Presentation(_io.BytesIO(file_bytes))
                    extra_notes = []
                    for i, slide in enumerate(prs.slides):
                        try:
                            notes = slide.notes_slide.notes_text_frame.text.strip()
                            if notes:
                                extra_notes.append(f"[Slide {i+1} Speaker Notes: {notes}]")
                        except Exception:
                            pass
                    if extra_notes:
                        content['text'] += '\n\n' + '\n'.join(extra_notes)
                    return content

                # ── Strategy 2: python-pptx text extraction + slide renders ──
                # Used when LibreOffice is unavailable
                logger.info("[PPTX] LibreOffice unavailable — using python-pptx extraction")
                prs = Presentation(_io.BytesIO(file_bytes))
                slides_text = []
                extracted_images = []

                for i, slide in enumerate(prs.slides):
                    slide_parts = [f"--- Slide {i+1} ---"]

                    # Extract ALL text from every shape type
                    for shape in slide.shapes:
                        # Text frames (titles, bullets, body)
                        if hasattr(shape, 'text_frame'):
                            for para in shape.text_frame.paragraphs:
                                para_text = para.text.strip()
                                if para_text:
                                    slide_parts.append(para_text)
                        elif hasattr(shape, 'text') and shape.text.strip():
                            slide_parts.append(shape.text.strip())

                        # Tables
                        if shape.shape_type == 19:  # TABLE
                            try:
                                table = shape.table
                                for row in table.rows:
                                    row_text = ' | '.join(
                                        cell.text.strip() for cell in row.cells
                                        if cell.text.strip()
                                    )
                                    if row_text:
                                        slide_parts.append(row_text)
                            except Exception:
                                pass

                        # Charts
                        if shape.shape_type == 3:  # CHART
                            try:
                                chart = shape.chart
                                if chart.has_title:
                                    slide_parts.append(f"[Chart: {chart.chart_title.text_frame.text}]")
                            except Exception:
                                pass

                        # Embedded pictures
                        if shape.shape_type == 13:  # PICTURE
                            try:
                                img_blob = shape.image.blob
                                img_ext = shape.image.ext
                                extracted_images.append({
                                    'data': img_blob,
                                    'page': i + 1,
                                    'ext': img_ext or 'png',
                                    'is_large': len(img_blob) > 50000,
                                })
                            except Exception:
                                pass

                    # Speaker notes
                    try:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_parts.append(f"[Speaker Notes: {notes_text}]")
                    except Exception:
                        pass

                    if len(slide_parts) > 1:
                        slides_text.append('\n'.join(slide_parts))

                content['text'] = '\n\n'.join(slides_text)
                content['page_count'] = len(prs.slides)

                # Render slides to images using Pillow fallback
                slide_renders = _render_slides_pillow(file_bytes, prs)
                if slide_renders:
                    content['slide_images'] = slide_renders
                elif extracted_images:
                    content['slide_images'] = extracted_images

            except ImportError:
                content['status'] = 'error'
                content['error'] = 'python-pptx not installed'
            except Exception as e:
                logger.error(f'[PPTX] Extraction error: {e}')
                content['status'] = 'error'
                content['error'] = str(e)
            
        else:
            content['status'] = 'error'
            content['error'] = f'Unsupported extension: {ext}'
            
    except Exception as e:
        logger.error(f'Extraction error for {extension}: {e}')
        content['status'] = 'error'
        content['error'] = str(e)
        
    return content


def _convert_pptx_to_pdf(file_bytes: bytes) -> bytes:
    """
    Convert PPTX to PDF using LibreOffice headless.
    Returns PDF bytes, or None if LibreOffice is unavailable.
    """
    import subprocess
    import tempfile
    import os

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, 'input.pptx')
            with open(pptx_path, 'wb') as f:
                f.write(file_bytes)

            result = subprocess.run(
                [
                    'libreoffice', '--headless', '--norestore',
                    '--convert-to', 'pdf',
                    '--outdir', tmpdir,
                    pptx_path
                ],
                capture_output=True,
                timeout=120,
            )

            if result.returncode == 0:
                pdf_path = os.path.join(tmpdir, 'input.pdf')
                if os.path.exists(pdf_path):
                    with open(pdf_path, 'rb') as f:
                        return f.read()
    except FileNotFoundError:
        # LibreOffice not installed
        pass
    except subprocess.TimeoutExpired:
        logger.warning("[PPTX→PDF] LibreOffice conversion timed out")
    except Exception as e:
        logger.warning(f"[PPTX→PDF] Conversion failed: {e}")

    return None


def _render_slides_pillow(file_bytes: bytes, prs) -> list:
    """
    Render PPTX slides as simple PNG images using Pillow.
    Creates a clean text layout image per slide.
    Used as fallback when LibreOffice is unavailable.
    """
    try:
        from PIL import Image, ImageDraw
        import io as _io

        renders = []
        for i, slide in enumerate(prs.slides):
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
                elif hasattr(shape, 'text') and shape.text.strip():
                    lines.append(shape.text.strip())

            img = Image.new('RGB', (1280, 720), color='white')
            draw = ImageDraw.Draw(img)
            y = 40
            for line in lines[:20]:
                draw.text((40, y), line[:100], fill='black')
                y += 34
                if y > 680:
                    break

            buf = _io.BytesIO()
            img.save(buf, format='PNG')
            renders.append({
                'data': buf.getvalue(),
                'page': i + 1,
                'ext': 'png',
                'is_large': True,
            })

        return renders
    except ImportError:
        return []
    except Exception as e:
        logger.warning(f"[PPTX Pillow render] Failed: {e}")
        return []
